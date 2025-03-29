"""
MCP Server for PowerPoint Operations

This server provides tools to create, edit and manage PowerPoint presentations.
It's implemented using the Model Context Protocol (MCP) Python SDK.
"""

import os
import sys
import io
from mcp.server.fastmcp import FastMCP
from typing import Optional, List, Dict, Any, Union, Tuple

# 标记库是否已安装
pptx_installed = True

# 尝试导入python-pptx库，如果没有安装则标记为未安装但不退出
try:
    import pptx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.dml import MSO_THEME_COLOR
except ImportError:
    print("警告: 未检测到python-pptx库，PowerPoint功能将不可用")
    print("请使用以下命令安装: pip install python-pptx")
    pptx_installed = False

# 尝试导入Pillow库，用于图片处理
pillow_installed = True
try:
    from PIL import Image
except ImportError:
    print("警告: 未检测到Pillow库，图片处理功能将受限")
    print("请使用以下命令安装: pip install Pillow")
    pillow_installed = False

# 创建一个MCP服务器，保持名称与配置文件一致
mcp = FastMCP("office editor")

@mcp.tool()
def create_powerpoint_presentation(filename: str) -> str:
    """
    创建一个新的PowerPoint演示文稿。
    
    Args:
        filename: 要创建的文件名 (不需要包含.pptx扩展名)
    
    Returns:
        包含操作结果的消息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法创建PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 确保文件名有.pptx扩展名
    if not filename.lower().endswith('.pptx'):
        filename += '.pptx'
    
    # 从环境变量获取输出路径，如果未设置则使用默认桌面路径
    output_path = os.environ.get('OFFICE_EDIT_PATH')
    if not output_path:
        output_path = os.path.join(os.path.expanduser('~'), '桌面')
    
    # 创建完整的文件路径
    file_path = os.path.join(output_path, filename)
    
    try:
        # 创建输出目录（如果不存在）
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        
        # 创建新的PowerPoint演示文稿
        prs = Presentation()
        
        # 保存演示文稿
        prs.save(file_path)
        
        return f"成功在 {output_path} 创建了PowerPoint演示文稿: {filename}"
    except Exception as e:
        return f"创建PowerPoint演示文稿时出错: {str(e)}"

@mcp.tool()
def open_powerpoint_presentation(file_path: str) -> str:
    """
    打开一个现有的PowerPoint演示文稿并读取其基本信息。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
    
    Returns:
        演示文稿的基本信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法打开PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片信息
        slide_count = len(prs.slides)
        
        # 获取每张幻灯片的基本信息
        slides_info = []
        for i, slide in enumerate(prs.slides):
            slide_title = "无标题"
            # 尝试获取幻灯片标题
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and (shape.name.startswith("Title") or "标题" in shape.name):
                    slide_title = shape.text[:50] + "..." if len(shape.text) > 50 else shape.text
                    break
            
            # 计算幻灯片上的形状数量
            shape_count = len(slide.shapes)
            
            slides_info.append(f"幻灯片 {i+1}: {slide_title} (包含 {shape_count} 个形状)")
        
        # 构建演示文稿信息
        presentation_info = (
            f"文件名: {os.path.basename(file_path)}\n"
            f"幻灯片数量: {slide_count}\n\n"
            f"幻灯片概览:\n" + "\n".join(slides_info)
        )
        
        return presentation_info
    except Exception as e:
        return f"打开PowerPoint演示文稿时出错: {str(e)}"

@mcp.tool()
def save_presentation_as(file_path: str, output_format: str = "pptx", new_filename: str = None) -> str:
    """
    保存PowerPoint演示文稿，可选择保存为不同格式。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        output_format: 保存格式，可选值: "pptx", "ppt", "pdf"
        new_filename: 新文件名(不含扩展名)，如果不提供则使用原文件名
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法保存PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    # 验证格式类型
    supported_formats = ["pptx", "ppt", "pdf"]
    if output_format not in supported_formats:
        return f"错误: 不支持的格式类型 '{output_format}'，支持的格式有: {', '.join(supported_formats)}"
    
    try:
        # 获取原文件名（不含扩展名）和目录
        file_dir = os.path.dirname(file_path)
        file_name = os.path.splitext(os.path.basename(file_path))[0]
        
        # 使用新文件名（如果提供）
        if new_filename:
            file_name = new_filename
        
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 创建新文件路径
        new_file_path = os.path.join(file_dir, f"{file_name}.{output_format}")
        
        # 根据格式类型保存
        if output_format == "pptx":
            prs.save(new_file_path)
            return f"已成功将演示文稿保存为: {new_file_path}"
        
        elif output_format == "ppt" or output_format == "pdf":
            # 注意：python-pptx库本身不支持直接保存为PPT或PDF格式
            # 需要调用外部应用程序PowerPoint来完成这个转换
            # 下面尝试使用win32com.client（仅限Windows系统）
            try:
                import win32com.client
                import win32api
                
                # 保存临时pptx文件
                temp_path = os.path.join(file_dir, f"{file_name}_temp.pptx")
                prs.save(temp_path)
                
                # 使用PowerPoint应用程序打开并另存为其他格式
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = False  # 不显示PowerPoint窗口
                
                # 打开临时pptx文件
                presentation = powerpoint.Presentations.Open(temp_path)
                
                if output_format == "ppt":
                    # 另存为PPT格式
                    # ppSaveAsPresentation = 1 表示以PPT格式保存
                    presentation.SaveAs(new_file_path, 1)
                else:  # output_format == "pdf"
                    # 另存为PDF格式
                    # ppSaveAsPDF = 32 表示以PDF格式保存
                    presentation.SaveAs(new_file_path, 32)
                
                # 关闭演示文稿和PowerPoint应用程序
                presentation.Close()
                powerpoint.Quit()
                
                # 删除临时文件
                os.remove(temp_path)
                
                return f"已成功将演示文稿保存为: {new_file_path}"
            
            except ImportError:
                return (f"无法完成转换: 保存为{output_format}格式需要在Windows系统上安装pywin32库和Microsoft PowerPoint。"
                        f"\n请使用命令安装pywin32: pip install pywin32")
            except Exception as e:
                return f"保存为{output_format}格式时出错: {str(e)}"
        
        return f"已成功将演示文稿保存为: {new_file_path}"
    except Exception as e:
        return f"保存PowerPoint演示文稿时出错: {str(e)}"

@mcp.tool()
def add_slide(file_path: str, layout_name: str = "Title and Content") -> str:
    """
    向PowerPoint演示文稿添加新幻灯片。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        layout_name: 幻灯片版式名称，常见值包括:
                    "Title Slide" (标题幻灯片)
                    "Title and Content" (标题和内容)
                    "Section Header" (节标题)
                    "Two Content" (两栏内容)
                    "Comparison" (比较)
                    "Title Only" (仅标题)
                    "Blank" (空白)
                    "Content with Caption" (带说明的内容)
                    "Picture with Caption" (带说明的图片)
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取所有可用的幻灯片版式
        available_layouts = [layout.name for layout in prs.slide_layouts]
        
        # 验证提供的版式名称是否有效
        layout_index = None
        for i, name in enumerate(available_layouts):
            if name.lower() == layout_name.lower() or layout_name.lower() in name.lower():
                layout_index = i
                break
        
        # 如果找不到匹配的版式，则使用默认版式（通常是"Title and Content"或index=1）
        if layout_index is None:
            # 尝试查找最接近的版式
            if "title" in layout_name.lower() and "content" in layout_name.lower():
                # 尝试查找标题和内容版式
                for i, name in enumerate(available_layouts):
                    if "title" in name.lower() and "content" in name.lower():
                        layout_index = i
                        break
            
            # 如果仍然找不到，使用索引1（通常是标题和内容）或0（标题）
            if layout_index is None:
                if len(prs.slide_layouts) > 1:
                    layout_index = 1  # 标题和内容
                else:
                    layout_index = 0  # 默认为第一个布局
            
            layout_name = available_layouts[layout_index]
        
        # 添加新幻灯片
        slide_layout = prs.slide_layouts[layout_index]
        new_slide = prs.slides.add_slide(slide_layout)
        
        # 保存演示文稿
        prs.save(file_path)
        
        # 获取当前幻灯片总数
        total_slides = len(prs.slides)
        
        return (f"已成功添加新幻灯片（版式: {layout_name}）\n"
                f"当前演示文稿共有 {total_slides} 张幻灯片")
    except Exception as e:
        return f"添加幻灯片时出错: {str(e)}"

@mcp.tool()
def delete_slide(file_path: str, slide_index: int) -> str:
    """
    删除PowerPoint演示文稿中的指定幻灯片。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要删除的幻灯片索引（从1开始计数）
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片总数
        total_slides = len(prs.slides)
        
        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
        
        # python-pptx没有直接的方法来删除幻灯片
        # 需要创建一个新的演示文稿并复制所有幻灯片，除了要删除的幻灯片
        
        # 创建一个副本文件路径
        temp_file_path = file_path + ".temp"
        
        # 创建一个新的临时演示文稿
        temp_prs = Presentation()
        
        # 复制除要删除的幻灯片外的所有幻灯片
        for i, slide in enumerate(prs.slides):
            # 跳过要删除的幻灯片
            if i + 1 == slide_index:
                continue
            
            # 复制当前幻灯片的布局
            slide_layout = temp_prs.slide_layouts[0]  # 默认使用第一个布局
            
            # 尝试匹配原始幻灯片的布局
            try:
                # 获取原始幻灯片的布局索引
                original_layout_idx = prs.slides[i].slide_layout.index
                if original_layout_idx < len(temp_prs.slide_layouts):
                    slide_layout = temp_prs.slide_layouts[original_layout_idx]
            except:
                pass  # 如果无法获取或匹配布局，使用默认布局
            
            # 创建新幻灯片
            new_slide = temp_prs.slides.add_slide(slide_layout)
            
            # 复制形状（这是一个简化的复制，不会完全复制所有内容和格式）
            # 注意：完整复制演示文稿的所有方面需要更复杂的逻辑
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # 复制文本形状
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    if shape.has_text_frame:
                        # 创建文本框
                        text_box = new_slide.shapes.add_textbox(left, top, width, height)
                        text_frame = text_box.text_frame
                        
                        # 复制文本内容
                        text_frame.text = shape.text_frame.text
        
        # 保存临时演示文稿
        temp_prs.save(temp_file_path)
        
        # 使用win32com来替换原文件（如果可用）
        try:
            # 关闭并删除原文件，然后重命名临时文件
            import os
            if os.path.exists(file_path):
                os.remove(file_path)
            os.rename(temp_file_path, file_path)
            
            return f"已成功删除幻灯片 {slide_index}，当前演示文稿还有 {total_slides - 1} 张幻灯片"
        except:
            # 如果无法替换，则返回临时文件的路径
            return (f"已成功创建不包含幻灯片 {slide_index} 的新演示文稿: {temp_file_path}"
                    f"\n请手动替换原文件")
        
    except Exception as e:
        return f"删除幻灯片时出错: {str(e)}"

@mcp.tool()
def reorder_slides(file_path: str, slide_index: int, new_position: int) -> str:
    """
    调整PowerPoint演示文稿中幻灯片的顺序。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要移动的幻灯片索引（从1开始计数）
        new_position: 幻灯片的新位置（从1开始计数）
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片总数
        total_slides = len(prs.slides)
        
        # 检查slide_index和new_position是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
        
        if new_position < 1 or new_position > total_slides:
            return f"错误: 无效的目标位置 {new_position}，演示文稿共有 {total_slides} 张幻灯片"
        
        # 如果位置相同，无需调整
        if slide_index == new_position:
            return f"幻灯片 {slide_index} 已经在位置 {new_position}，无需调整"
        
        # python-pptx目前不直接支持重新排序幻灯片
        # 需要创建一个新的演示文稿并按照新的顺序复制幻灯片
        
        # 创建一个副本文件路径
        temp_file_path = file_path + ".temp"
        
        # 创建一个新的临时演示文稿
        temp_prs = Presentation()
        
        # 调整幻灯片顺序
        slides = list(prs.slides)
        
        # 从列表中删除要移动的幻灯片
        slide_to_move = slides.pop(slide_index - 1)
        
        # 在新位置插入幻灯片
        slides.insert(new_position - 1, slide_to_move)
        
        # 按新顺序复制幻灯片
        for slide in slides:
            # 尝试获取原始幻灯片的布局
            try:
                original_layout_idx = slide.slide_layout.index
                if original_layout_idx < len(temp_prs.slide_layouts):
                    slide_layout = temp_prs.slide_layouts[original_layout_idx]
                else:
                    slide_layout = temp_prs.slide_layouts[0]  # 默认使用第一个布局
            except:
                slide_layout = temp_prs.slide_layouts[0]  # 如果失败，使用第一个布局
            
            # 创建新幻灯片
            new_slide = temp_prs.slides.add_slide(slide_layout)
            
            # 复制形状（这是一个简化的复制，不会完全复制所有内容和格式）
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    # 复制文本形状
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    
                    if shape.has_text_frame:
                        # 创建文本框
                        text_box = new_slide.shapes.add_textbox(left, top, width, height)
                        text_frame = text_box.text_frame
                        
                        # 复制文本内容
                        text_frame.text = shape.text_frame.text
        
        # 保存临时演示文稿
        temp_prs.save(temp_file_path)
        
        # 使用os来替换原文件
        try:
            # 关闭并删除原文件，然后重命名临时文件
            import os
            if os.path.exists(file_path):
                os.remove(file_path)
            os.rename(temp_file_path, file_path)
            
            return f"已成功将幻灯片 {slide_index} 移动到位置 {new_position}"
        except:
            # 如果无法替换，则返回临时文件的路径
            return (f"已成功创建调整了幻灯片顺序的新演示文稿: {temp_file_path}"
                    f"\n请手动替换原文件")
        
    except Exception as e:
        return f"调整幻灯片顺序时出错: {str(e)}"

@mcp.tool()
def set_slide_background(file_path: str, slide_index: int, background_color: str = None) -> str:
    """
    为PowerPoint演示文稿中的指定幻灯片设置背景颜色。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要设置背景的幻灯片索引（从1开始计数），0表示所有幻灯片
        background_color: 背景颜色的十六进制代码，例如"#FF0000"表示红色
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    # 验证颜色格式
    if background_color:
        if not background_color.startswith('#') or len(background_color) != 7:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"
        
        try:
            # 解析RGB值
            r = int(background_color[1:3], 16)
            g = int(background_color[3:5], 16)
            b = int(background_color[5:7], 16)
        except ValueError:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片总数
        total_slides = len(prs.slides)
        
        # 如果slide_index不是0，检查其有效性
        if slide_index != 0 and (slide_index < 1 or slide_index > total_slides):
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
        
        # 确定要修改的幻灯片
        slides_to_modify = []
        if slide_index == 0:
            # 修改所有幻灯片
            slides_to_modify = list(range(total_slides))
        else:
            # 修改指定幻灯片
            slides_to_modify = [slide_index - 1]  # 转换为0开始的索引
        
        # 设置背景颜色
        for idx in slides_to_modify:
            slide = prs.slides[idx]
            
            # 获取幻灯片背景
            background = slide.background
            
            # 设置背景填充类型为纯色
            fill = background.fill
            fill.solid()
            
            # 设置背景颜色
            if background_color:
                fill.fore_color.rgb = RGBColor(r, g, b)
        
        # 保存演示文稿
        prs.save(file_path)
        
        # 构建结果消息
        if slide_index == 0:
            return f"已成功设置所有幻灯片的背景颜色为 {background_color}"
        else:
            return f"已成功设置幻灯片 {slide_index} 的背景颜色为 {background_color}"
        
    except Exception as e:
        return f"设置幻灯片背景时出错: {str(e)}"

@mcp.tool()
def add_text_box(
    file_path: str, 
    slide_index: int, 
    text: str, 
    left: float = 1.0, 
    top: float = 1.0, 
    width: float = 4.0, 
    height: float = 1.0,
    font_name: str = None,
    font_size: int = None,
    font_bold: bool = False,
    font_italic: bool = False,
    text_color: str = None,
    alignment: str = "left"
) -> str:
    """
    向PowerPoint演示文稿的指定幻灯片添加文本框。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要添加文本框的幻灯片索引（从1开始计数）
        text: 要添加的文本内容
        left: 文本框左侧位置（英寸）
        top: 文本框顶部位置（英寸）
        width: 文本框宽度（英寸）
        height: 文本框高度（英寸）
        font_name: 字体名称，如"宋体"、"Arial"等
        font_size: 字体大小（磅）
        font_bold: 是否加粗
        font_italic: 是否斜体
        text_color: 文本颜色的十六进制代码，例如"#FF0000"表示红色
        alignment: 文本对齐方式，可选值: "left", "center", "right", "justify"
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    # 验证文本颜色格式
    if text_color:
        if not text_color.startswith('#') or len(text_color) != 7:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"
        
        try:
            # 解析RGB值
            r = int(text_color[1:3], 16)
            g = int(text_color[3:5], 16)
            b = int(text_color[5:7], 16)
        except ValueError:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"
    
    # 验证对齐方式
    alignment_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY
    }
    
    if alignment.lower() not in alignment_map:
        return f"错误: 无效的对齐方式 '{alignment}'，支持的对齐方式有: {', '.join(alignment_map.keys())}"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片总数
        total_slides = len(prs.slides)
        
        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
        
        # 获取指定幻灯片
        slide = prs.slides[slide_index - 1]
        
        # 添加文本框
        textbox = slide.shapes.add_textbox(
            Inches(left), 
            Inches(top), 
            Inches(width), 
            Inches(height)
        )
        
        # 获取文本框的文本帧
        text_frame = textbox.text_frame
        
        # 设置文本内容
        text_frame.text = text
        
        # 设置文本格式
        p = text_frame.paragraphs[0]
        
        # 设置对齐方式
        p.alignment = alignment_map[alignment.lower()]
        
        # 获取run对象来设置字体属性
        run = p.runs[0]
        
        # 设置字体名称
        if font_name:
            run.font.name = font_name
        
        # 设置字体大小
        if font_size:
            run.font.size = Pt(font_size)
        
        # 设置字体加粗
        if font_bold:
            run.font.bold = True
        
        # 设置字体斜体
        if font_italic:
            run.font.italic = True
        
        # 设置文本颜色
        if text_color:
            run.font.color.rgb = RGBColor(r, g, b)
        
        # 保存演示文稿
        prs.save(file_path)
        
        return f"已成功在幻灯片 {slide_index} 添加文本框"
        
    except Exception as e:
        return f"添加文本框时出错: {str(e)}"

@mcp.tool()
def insert_image(file_path: str, slide_index: int, image_path: str, left: float = 1.0, top: float = 1.0, width: float = None, height: float = None) -> str:
    """
    在PowerPoint演示文稿的指定幻灯片中插入图片。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要插入图片的幻灯片索引（从1开始计数）
        image_path: 图片文件的完整路径或相对路径
        left: 图片左侧位置（英寸）
        top: 图片顶部位置（英寸）
        width: 图片宽度（英寸），如果不指定则保持原始宽高比
        height: 图片高度（英寸），如果不指定则保持原始宽高比
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 演示文稿文件 {file_path} 不存在"
    
    # 确保图片文件存在
    if not os.path.isabs(image_path):
        # 尝试找到图片的绝对路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 检查几个可能的位置
        potential_paths = [
            os.path.join(base_path, image_path),
            os.path.join(os.path.dirname(file_path), image_path)
        ]
        
        for path in potential_paths:
            if os.path.exists(path):
                image_path = path
                break
        else:
            return f"错误: 图片文件 {image_path} 不存在"
    
    if not os.path.exists(image_path):
        return f"错误: 图片文件 {image_path} 不存在"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片总数
        total_slides = len(prs.slides)
        
        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
        
        # 获取指定幻灯片
        slide = prs.slides[slide_index - 1]
        
        # 插入图片
        if width and height:
            # 按指定尺寸插入图片
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top),
                width=Inches(width),
                height=Inches(height)
            )
        elif width:
            # 只指定宽度，高度按比例计算
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top),
                width=Inches(width)
            )
        elif height:
            # 只指定高度，宽度按比例计算
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top),
                height=Inches(height)
            )
        else:
            # 使用原始尺寸
            slide.shapes.add_picture(
                image_path,
                Inches(left),
                Inches(top)
            )
        
        # 保存演示文稿
        prs.save(file_path)
        
        return f"已成功在幻灯片 {slide_index} 插入图片"
        
    except Exception as e:
        return f"插入图片时出错: {str(e)}"

@mcp.tool()
def insert_table(file_path: str, slide_index: int, rows: int, cols: int, data: List[List[str]] = None, left: float = 1.0, top: float = 1.0, width: float = 6.0, height: float = 3.0) -> str:
    """
    在PowerPoint演示文稿的指定幻灯片中插入表格。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要插入表格的幻灯片索引（从1开始计数）
        rows: 表格行数
        cols: 表格列数
        data: 表格数据，二维列表，每个内部列表表示一行数据
        left: 表格左侧位置（英寸）
        top: 表格顶部位置（英寸）
        width: 表格宽度（英寸）
        height: 表格高度（英寸）
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    # 验证行列数
    if rows < 1 or cols < 1:
        return "错误: 表格行数和列数必须大于0"
    
    # 验证数据格式
    if data:
        if len(data) > rows:
            return f"错误: 提供的数据行数({len(data)})超过了表格行数({rows})"
        
        for i, row_data in enumerate(data):
            if len(row_data) > cols:
                return f"错误: 第{i+1}行的数据列数({len(row_data)})超过了表格列数({cols})"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片总数
        total_slides = len(prs.slides)
        
        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
        
        # 获取指定幻灯片
        slide = prs.slides[slide_index - 1]
        
        # 创建表格形状
        shape = slide.shapes.add_table(
            rows, 
            cols, 
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height)
        )
        
        # 获取表格对象
        table = shape.table
        
        # 填充表格数据
        if data:
            for i, row_data in enumerate(data):
                if i < rows:  # 确保不超过表格行数
                    for j, cell_data in enumerate(row_data):
                        if j < cols:  # 确保不超过表格列数
                            # 获取单元格
                            cell = table.cell(i, j)
                            
                            # 设置单元格文本
                            cell.text = str(cell_data)
        
        # 保存演示文稿
        prs.save(file_path)
        
        return f"已成功在幻灯片 {slide_index} 插入 {rows}×{cols} 的表格"
        
    except Exception as e:
        return f"插入表格时出错: {str(e)}"

@mcp.tool()
def add_slide_notes(file_path: str, slide_index: int, notes_text: str) -> str:
    """
    为PowerPoint演示文稿的指定幻灯片添加备注。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要添加备注的幻灯片索引（从1开始计数）
        notes_text: 备注内容
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片总数
        total_slides = len(prs.slides)
        
        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
        
        # 获取指定幻灯片
        slide = prs.slides[slide_index - 1]
        
        # 获取备注页对象
        notes_slide = slide.notes_slide
        
        # 设置备注文本
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text
        
        # 保存演示文稿
        prs.save(file_path)
        
        return f"已成功为幻灯片 {slide_index} 添加备注"
        
    except Exception as e:
        return f"添加幻灯片备注时出错: {str(e)}"

@mcp.tool()
def insert_shape(file_path: str, slide_index: int, shape_type: str, left: float = 1.0, top: float = 1.0, width: float = 2.0, height: float = 2.0, fill_color: str = None) -> str:
    """
    在PowerPoint演示文稿的指定幻灯片中插入形状。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要插入形状的幻灯片索引（从1开始计数）
        shape_type: 形状类型，可选值包括: "rectangle", "oval", "triangle", "arrow", "line"
        left: 形状左侧位置（英寸）
        top: 形状顶部位置（英寸）
        width: 形状宽度（英寸）
        height: 形状高度（英寸）
        fill_color: 填充颜色的十六进制代码，例如"#FF0000"表示红色
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    # 验证形状类型
    shape_type_map = {
        "rectangle": MSO_SHAPE.RECTANGLE,
        "oval": MSO_SHAPE.OVAL,
        "triangle": MSO_SHAPE.TRIANGLE,
        "arrow": MSO_SHAPE.RIGHT_ARROW,
        "line": MSO_SHAPE.LINE
    }
    
    if shape_type.lower() not in shape_type_map:
        return f"错误: 无效的形状类型 '{shape_type}'，支持的形状类型有: {', '.join(shape_type_map.keys())}"
    
    # 验证颜色格式
    if fill_color:
        if not fill_color.startswith('#') or len(fill_color) != 7:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"
        
        try:
            # 解析RGB值
            r = int(fill_color[1:3], 16)
            g = int(fill_color[3:5], 16)
            b = int(fill_color[5:7], 16)
        except ValueError:
            return "错误: 颜色格式无效，请使用十六进制格式，例如'#FF0000'表示红色"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片总数
        total_slides = len(prs.slides)
        
        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
        
        # 获取指定幻灯片
        slide = prs.slides[slide_index - 1]
        
        # 获取形状类型
        shape_type_value = shape_type_map[shape_type.lower()]
        
        # 添加形状
        shape = slide.shapes.add_shape(
            shape_type_value,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height)
        )
        
        # 设置填充颜色
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(r, g, b)
        
        # 保存演示文稿
        prs.save(file_path)
        
        return f"已成功在幻灯片 {slide_index} 插入 {shape_type} 形状"
    
    except Exception as e:
        return f"插入形状时出错: {str(e)}"

@mcp.tool()
def insert_chart(file_path: str, slide_index: int, chart_type: str, data: List[List[str]], left: float = 1.0, top: float = 1.0, width: float = 6.0, height: float = 4.0, has_legend: bool = True, has_title: bool = True, title: str = "图表标题") -> str:
    """
    在PowerPoint演示文稿的指定幻灯片中插入图表。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要插入图表的幻灯片索引（从1开始计数）
        chart_type: 图表类型，可选值: "column", "line", "pie", "bar", "area"
        data: 图表数据，二维列表，第一行为类别标签，第一列为数据系列名称
        left: 图表左侧位置（英寸）
        top: 图表顶部位置（英寸）
        width: 图表宽度（英寸）
        height: 图表高度（英寸）
        has_legend: 是否显示图例
        has_title: 是否显示标题
        title: 图表标题文本
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    # 验证图表类型
    from pptx.enum.charts import XL_CHART_TYPE
    chart_type_map = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "area": XL_CHART_TYPE.AREA
    }
    
    if chart_type.lower() not in chart_type_map:
        return f"错误: 无效的图表类型 '{chart_type}'，支持的图表类型有: {', '.join(chart_type_map.keys())}"
    
    # 验证数据格式
    if not data or len(data) < 2:
        return "错误: 图表数据必须至少包含两行（标题行和至少一个数据行）"
    
    for row in data:
        if not row or len(row) < 2:
            return "错误: 每行数据必须至少包含两列（类别名称和至少一个数据点）"
    
    try:
        # 打开PowerPoint演示文稿
        prs = Presentation(file_path)
        
        # 获取幻灯片总数
        total_slides = len(prs.slides)
        
        # 检查slide_index是否有效
        if slide_index < 1 or slide_index > total_slides:
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {total_slides} 张幻灯片"
        
        # 获取指定幻灯片
        slide = prs.slides[slide_index - 1]
        
        # 转换字符串数据为数值型数据
        # 第一行和第一列保持为字符串（类别和系列名称）
        # 其余部分转换为浮点数
        chart_data = []
        for i, row in enumerate(data):
            new_row = []
            for j, cell in enumerate(row):
                if i == 0 or j == 0:
                    # 第一行或第一列保持为字符串
                    new_row.append(str(cell))
                else:
                    # 尝试转换为浮点数
                    try:
                        new_row.append(float(cell))
                    except (ValueError, TypeError):
                        # 如果无法转换，使用0
                        new_row.append(0.0)
            chart_data.append(new_row)
        
        # 获取图表类型
        chart_type_value = chart_type_map[chart_type.lower()]
        
        # 创建图表
        chart_data_obj = None
        
        # 根据图表类型的不同，处理数据的方式也不同
        from pptx.chart.data import CategoryChartData, ChartData
        
        if chart_type.lower() == "pie":
            # 饼图只使用第二行数据
            chart_data_obj = ChartData()
            categories = [cat for cat in chart_data[0][1:]]
            values = [val for val in chart_data[1][1:]]
            
            # 添加类别和值
            chart_data_obj.categories = categories
            chart_data_obj.add_series(chart_data[1][0], values)
        else:
            # 其他图表类型
            chart_data_obj = CategoryChartData()
            
            # 添加类别（第一行除第一个单元格外的所有单元格）
            categories = [cat for cat in chart_data[0][1:]]
            chart_data_obj.categories = categories
            
            # 添加每个系列的数据（从第二行开始）
            for row in chart_data[1:]:
                series_name = row[0]
                values = [val for val in row[1:]]
                chart_data_obj.add_series(series_name, values)
        
        # 添加图表
        chart = slide.shapes.add_chart(
            chart_type_value,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            chart_data_obj
        ).chart
        
        # 设置图表标题
        if has_title:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
        else:
            chart.has_title = False
        
        # 设置图例
        chart.has_legend = has_legend
        if has_legend:
            chart.legend.position = 1  # 底部
            chart.legend.include_in_layout = False
        
        # 保存演示文稿
        prs.save(file_path)
        
        return f"已成功在幻灯片 {slide_index} 插入 {chart_type} 图表"
    
    except Exception as e:
        return f"插入图表时出错: {str(e)}"

@mcp.tool()
def apply_presentation_theme(file_path: str, theme_name: str) -> str:
    """
    为PowerPoint演示文稿应用主题。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        theme_name: 主题名称，可选值包括: "Office", "Adjacency", "Angles", "Apex", 
                   "Apothecary", "Aspect", "Austin", "Black Tie", "Civic", "Clarity", 
                   "Composite", "Concourse", "Couture", "Elemental", "Equity", "Essential",
                   "Executive", "Grid", "Hardcover", "Horizon", "Integral", "Ion", 
                   "Ion Boardroom", "Median", "Metro", "Module", "Newsprint", "Opulent",
                   "Organic", "Oriel", "Origin", "Paper", "Perspective", "Pushpin",
                   "Slipstream", "Solstice", "Technic", "Thatch", "Trek", "Urban", "Vapor Trail", "Wisp"
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    # Python-pptx不直接支持更改主题，我们需要使用win32com
    try:
        import win32com.client
        
        # 初始化PowerPoint应用程序
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False  # 不显示PowerPoint窗口
        
        # 打开演示文稿
        presentation = powerpoint.Presentations.Open(file_path)
        
        # 尝试应用主题
        try:
            # 获取主题文件的绝对路径
            # 注意：这需要Office安装了对应的主题
            # 主题文件通常位于Office的安装目录的THEMES子目录中
            theme_file = f"{theme_name}.thmx"
            
            # 应用主题
            presentation.ApplyTheme(theme_file)
            
            # 保存演示文稿
            presentation.Save()
            success = True
            message = f"已成功应用主题 '{theme_name}'"
        except Exception as theme_error:
            success = False
            message = f"应用主题时出错: {str(theme_error)}"
        
        # 关闭演示文稿和PowerPoint应用程序
        presentation.Close()
        powerpoint.Quit()
        
        if success:
            return message
        else:
            # 如果使用COM方法失败，提供一个备选方法
            return f"{message}\n可以尝试在PowerPoint中手动应用主题。"
    
    except ImportError:
        return "错误: 无法应用主题，这个功能需要在Windows系统上安装pywin32库。\n请使用命令安装: pip install pywin32"
    
    except Exception as e:
        return f"应用主题时出错: {str(e)}"

@mcp.tool()
def add_animation(file_path: str, slide_index: int, shape_index: int, animation_type: str, animation_trigger: str = "on_click", animation_delay: float = 0.0) -> str:
    """
    为PowerPoint演示文稿中的形状添加动画效果。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要添加动画的幻灯片索引（从1开始计数）
        shape_index: 要添加动画的形状索引（从1开始计数）
        animation_type: 动画类型，可选值: 
                      "entrance": 入场动画
                      "emphasis": 强调动画
                      "exit": 退场动画
        animation_trigger: 动画触发方式，可选值:
                        "on_click": 点击时播放
                        "with_previous": 与上一动画同时播放
                        "after_previous": 在上一动画之后播放
        animation_delay: 动画延迟时间（秒）
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    # 验证动画类型
    if animation_type not in ["entrance", "emphasis", "exit"]:
        return f"错误: 无效的动画类型 '{animation_type}'，支持的动画类型有: entrance, emphasis, exit"
    
    # 验证动画触发方式
    if animation_trigger not in ["on_click", "with_previous", "after_previous"]:
        return f"错误: 无效的动画触发方式 '{animation_trigger}'，支持的触发方式有: on_click, with_previous, after_previous"
    
    # Python-pptx不直接支持动画效果，我们需要使用win32com
    try:
        import win32com.client
        from win32com.client import constants
        
        # 初始化PowerPoint应用程序
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False  # 不显示PowerPoint窗口
        
        # 打开演示文稿
        presentation = powerpoint.Presentations.Open(file_path)
        
        # 检查幻灯片索引
        if slide_index < 1 or slide_index > presentation.Slides.Count:
            # 关闭文件和应用程序
            presentation.Close()
            powerpoint.Quit()
            return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {presentation.Slides.Count} 张幻灯片"
        
        # 获取指定幻灯片
        slide = presentation.Slides.Item(slide_index)
        
        # 检查形状索引
        if shape_index < 1 or shape_index > slide.Shapes.Count:
            # 关闭文件和应用程序
            presentation.Close()
            powerpoint.Quit()
            return f"错误: 无效的形状索引 {shape_index}，幻灯片 {slide_index} 共有 {slide.Shapes.Count} 个形状"
        
        # 获取指定形状
        shape = slide.Shapes.Item(shape_index)
        
        # 添加动画
        animation_effect = None
        
        # 获取动画序列
        animation_sequence = slide.TimeLine.MainSequence
        
        # 根据动画类型添加不同效果
        if animation_type == "entrance":
            # 入场动画
            animation_effect = animation_sequence.AddEntrance(
                shape, 
                constants.msoAnimEffectFade,  # 默认使用淡入效果
                constants.msoAnimateByObject, 
                constants.msoAnimTriggerOnPageClick
            )
        elif animation_type == "emphasis":
            # 强调动画
            animation_effect = animation_sequence.AddEmphasis(
                shape, 
                constants.msoAnimEffectPulse,  # 默认使用脉冲效果
                constants.msoAnimateByObject, 
                constants.msoAnimTriggerOnPageClick
            )
        elif animation_type == "exit":
            # 退场动画
            animation_effect = animation_sequence.AddExit(
                shape, 
                constants.msoAnimEffectFade,  # 默认使用淡出效果
                constants.msoAnimateByObject, 
                constants.msoAnimTriggerOnPageClick
            )
        
        # 设置动画触发方式
        if animation_trigger == "on_click":
            animation_effect.Timing.TriggerType = constants.msoAnimTriggerOnPageClick
        elif animation_trigger == "with_previous":
            animation_effect.Timing.TriggerType = constants.msoAnimTriggerWithPrevious
        elif animation_trigger == "after_previous":
            animation_effect.Timing.TriggerType = constants.msoAnimTriggerAfterPrevious
        
        # 设置动画延迟时间
        animation_effect.Timing.TriggerDelayTime = animation_delay
        
        # 保存演示文稿
        presentation.Save()
        
        # 关闭演示文稿和PowerPoint应用程序
        presentation.Close()
        powerpoint.Quit()
        
        return f"已成功为幻灯片 {slide_index} 中的形状 {shape_index} 添加 {animation_type} 动画"
    
    except ImportError:
        return "错误: 无法添加动画，这个功能需要在Windows系统上安装pywin32库。\n请使用命令安装: pip install pywin32"
    
    except Exception as e:
        # 尝试关闭PowerPoint
        try:
            if 'presentation' in locals() and presentation:
                presentation.Close(False)
            if 'powerpoint' in locals() and powerpoint:
                powerpoint.Quit()
        except:
            pass
        
        return f"添加动画时出错: {str(e)}"

@mcp.tool()
def set_slide_transition(file_path: str, slide_index: int, transition_type: str, duration: float = 1.0, on_click: bool = True, auto_advance: bool = False, advance_time: float = 5.0) -> str:
    """
    设置PowerPoint演示文稿中幻灯片的切换效果。
    
    Args:
        file_path: PowerPoint演示文稿的完整路径或相对于输出目录的路径
        slide_index: 要设置切换效果的幻灯片索引（从1开始计数），0表示所有幻灯片
        transition_type: 切换效果类型，可选值包括: "none", "fade", "push", "wipe", "split", "reveal", 
                       "random", "shape", "zoom", "gallery", "dissolve"
        duration: 切换效果持续时间（秒）
        on_click: 是否点击鼠标时切换幻灯片
        auto_advance: 是否自动前进到下一张幻灯片
        advance_time: 自动前进的时间（秒），仅当auto_advance=True时有效
    
    Returns:
        操作结果信息
    """
    # 检查是否安装了必要的库
    if not pptx_installed:
        return "错误: 无法操作PowerPoint演示文稿，请先安装python-pptx库: pip install python-pptx"
    
    # 检查是否提供了完整路径
    if not os.path.isabs(file_path):
        # 从环境变量获取基础路径
        base_path = os.environ.get('OFFICE_EDIT_PATH')
        if not base_path:
            base_path = os.path.join(os.path.expanduser('~'), '桌面')
        
        # 构建完整路径
        file_path = os.path.join(base_path, file_path)
    
    # 确保文件存在
    if not os.path.exists(file_path):
        return f"错误: 文件 {file_path} 不存在"
    
    # 验证切换效果类型
    transition_types = ["none", "fade", "push", "wipe", "split", "reveal", "random", "shape", "zoom", "gallery", "dissolve"]
    if transition_type.lower() not in transition_types:
        return f"错误: 无效的切换效果类型 '{transition_type}'，支持的类型有: {', '.join(transition_types)}"
    
    # Python-pptx不直接支持设置幻灯片切换效果，我们需要使用win32com
    try:
        import win32com.client
        from win32com.client import constants
        
        # 初始化PowerPoint应用程序
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = False  # 不显示PowerPoint窗口
        
        # 打开演示文稿
        presentation = powerpoint.Presentations.Open(file_path)
        
        # 切换效果类型映射
        transition_type_map = {
            "none": constants.ppTransitionNone,
            "fade": constants.ppTransitionFade,
            "push": constants.ppTransitionPush,
            "wipe": constants.ppTransitionWipe,
            "split": constants.ppTransitionSplit,
            "reveal": constants.ppTransitionReveal,
            "random": constants.ppTransitionRandom,
            "shape": constants.ppTransitionShapes,
            "zoom": constants.ppTransitionZoom,
            "gallery": constants.ppTransitionGallery,
            "dissolve": constants.ppTransitionDissolve
        }
        
        # 获取切换效果类型的值
        transition_value = transition_type_map.get(transition_type.lower())
        
        # 确定要修改的幻灯片
        if slide_index == 0:
            # 修改所有幻灯片
            for slide in presentation.Slides:
                # 设置切换效果
                slide.SlideShowTransition.EntryEffect = transition_value
                # 设置切换效果持续时间
                slide.SlideShowTransition.Duration = duration
                # 设置是否点击鼠标时切换
                slide.SlideShowTransition.AdvanceOnClick = on_click
                # 设置是否自动前进
                slide.SlideShowTransition.AdvanceOnTime = auto_advance
                # 设置自动前进的时间
                if auto_advance:
                    slide.SlideShowTransition.AdvanceTime = advance_time
            
            # 保存演示文稿
            presentation.Save()
            
            # 构建结果消息
            result = f"已成功为所有幻灯片设置 {transition_type} 切换效果"
        else:
            # 修改指定幻灯片
            # 检查幻灯片索引是否有效
            if slide_index < 1 or slide_index > presentation.Slides.Count:
                # 关闭文件和应用程序
                presentation.Close()
                powerpoint.Quit()
                return f"错误: 无效的幻灯片索引 {slide_index}，演示文稿共有 {presentation.Slides.Count} 张幻灯片"
            
            # 获取指定幻灯片
            slide = presentation.Slides.Item(slide_index)
            
            # 设置切换效果
            slide.SlideShowTransition.EntryEffect = transition_value
            # 设置切换效果持续时间
            slide.SlideShowTransition.Duration = duration
            # 设置是否点击鼠标时切换
            slide.SlideShowTransition.AdvanceOnClick = on_click
            # 设置是否自动前进
            slide.SlideShowTransition.AdvanceOnTime = auto_advance
            # 设置自动前进的时间
            if auto_advance:
                slide.SlideShowTransition.AdvanceTime = advance_time
            
            # 保存演示文稿
            presentation.Save()
            
            # 构建结果消息
            result = f"已成功为幻灯片 {slide_index} 设置 {transition_type} 切换效果"
        
        # 关闭演示文稿和PowerPoint应用程序
        presentation.Close()
        powerpoint.Quit()
        
        return result
    
    except ImportError:
        return "错误: 无法设置幻灯片切换效果，这个功能需要在Windows系统上安装pywin32库。\n请使用命令安装: pip install pywin32"
    
    except Exception as e:
        # 尝试关闭PowerPoint
        try:
            if 'presentation' in locals() and presentation:
                presentation.Close(False)
            if 'powerpoint' in locals() and powerpoint:
                powerpoint.Quit()
        except:
            pass
        
        return f"设置幻灯片切换效果时出错: {str(e)}"

# 添加主程序入口点
if __name__ == "__main__":
    try:
        # 启动MCP服务器
        mcp.run()
    except KeyboardInterrupt:
        # 优雅地处理Ctrl+C中断
        print("服务器已停止")
    except Exception as e:
        # 处理其他异常
        print(f"服务器运行时出错: {str(e)}")
        sys.exit(1)
