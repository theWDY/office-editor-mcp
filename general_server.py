"""
MCP Server for General Office Operations

This server provides advanced tools for Office document processing including OCR, 
document comparison, translation, encryption, and database operations.
It's implemented using the Model Context Protocol (MCP) Python SDK.
"""

import os
import sys
import json
import logging
import time
import base64
import io
from mcp.server.fastmcp import FastMCP
from typing import Dict, Any, List, Union, Tuple

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[logging.StreamHandler()]
)

logger = logging.getLogger("GeneralServer")

# 导入Office操作相关库
try:
    import docx
    from docx.shared import Pt, RGBColor, Inches
    import openpyxl
    from openpyxl.utils import get_column_letter
    from openpyxl.styles import Font, PatternFill, Alignment
    from pptx import Presentation
    from pptx.util import Inches as PptxInches
    import pdf2docx
    import PyPDF2
    import pytesseract
    from PIL import Image
    import difflib
    import deepl
    import cryptography
    from cryptography.fernet import Fernet
    import sqlalchemy
    from sqlalchemy import create_engine, text, MetaData, Table, Column, inspect
    import threading
    import concurrent.futures
    import pandas as pd
    import shutil
    import hashlib
except ImportError as e:
    logger.error(f"导入错误: {e}")
    logger.info("请安装所需依赖: pip install python-docx openpyxl python-pptx pdf2docx PyPDF2 pytesseract Pillow difflib deepl cryptography sqlalchemy pandas")

# 全局变量
OUTPUT_DIR = os.environ.get("OFFICE_EDIT_PATH", os.path.expanduser("~"))
logger.info(f"输出目录设置为: {OUTPUT_DIR}")

# 创建一个MCP服务器
mcp = FastMCP("office-editor")

# 工具函数
def extract_document_text(doc_path: str) -> str:
    """从不同类型的文档中提取文本内容"""
    ext = os.path.splitext(doc_path)[1].lower()
    
    try:
        if ext in ['.txt', '.md', '.json', '.py', '.xml', '.html', '.css', '.js']:
            # 纯文本文件
            with open(doc_path, 'r', encoding='utf-8') as f:
                return f.read()
                
        elif ext in ['.docx']:
            # Word文档
            doc = docx.Document(doc_path)
            return '\n'.join([para.text for para in doc.paragraphs])
            
        elif ext in ['.xlsx', '.xls']:
            # Excel文档
            wb = openpyxl.load_workbook(doc_path, data_only=True)
            text = []
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text.append(f"Sheet: {sheet_name}")
                for row in sheet.iter_rows():
                    text.append('\t'.join([str(cell.value) if cell.value is not None else '' for cell in row]))
            return '\n'.join(text)
            
        elif ext in ['.pptx', '.ppt']:
            # PowerPoint文档
            prs = Presentation(doc_path)
            text = []
            for i, slide in enumerate(prs.slides):
                text.append(f"Slide {i+1}")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return '\n'.join(text)
            
        elif ext in ['.pdf']:
            # PDF文档
            text = []
            with open(doc_path, 'rb') as file:
                reader = PyPDF2.PdfReader(file)
                for page_num in range(len(reader.pages)):
                    page = reader.pages[page_num]
                    text.append(page.extract_text())
            return '\n'.join(text)
            
        else:
            logger.warning(f"不支持的文档类型: {ext}")
            return None
            
    except Exception as e:
        logger.error(f"提取文档文本时出错: {str(e)}")
        return None

def replace_placeholders(file_path: str, data_mapping: Dict[str, List[str]], index: int):
    """替换文档中的占位符"""
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == '.docx':
            # Word文档占位符替换
            doc = docx.Document(file_path)
            
            for paragraph in doc.paragraphs:
                for key, values in data_mapping.items():
                    placeholder = f"{{{key}}}"
                    if placeholder in paragraph.text and index < len(values):
                        paragraph.text = paragraph.text.replace(placeholder, values[index])
            
            # 处理表格中的占位符
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.paragraphs:
                            for key, values in data_mapping.items():
                                placeholder = f"{{{key}}}"
                                if placeholder in paragraph.text and index < len(values):
                                    paragraph.text = paragraph.text.replace(placeholder, values[index])
            
            doc.save(file_path)
            
        elif ext in ['.xlsx', '.xls']:
            # Excel文档占位符替换
            wb = openpyxl.load_workbook(file_path)
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row in sheet.iter_rows():
                    for cell in row:
                        if cell.value and isinstance(cell.value, str):
                            for key, values in data_mapping.items():
                                placeholder = f"{{{key}}}"
                                if placeholder in cell.value and index < len(values):
                                    cell.value = cell.value.replace(placeholder, values[index])
            
            wb.save(file_path)
            
        elif ext in ['.pptx', '.ppt']:
            # PowerPoint文档占位符替换
            prs = Presentation(file_path)
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for key, values in data_mapping.items():
                                    placeholder = f"{{{key}}}"
                                    if placeholder in run.text and index < len(values):
                                        run.text = run.text.replace(placeholder, values[index])
            
            prs.save(file_path)
            
    except Exception as e:
        logger.error(f"替换占位符过程中出错: {str(e)}")
        raise

# OCR识别功能
@mcp.tool()
def ocr_recognize_text(image_path: str, language: str = "chi_sim+eng") -> Dict[str, Any]:
    """
    从图片中提取文本内容
    
    Args:
        image_path: 图片文件路径
        language: OCR识别语言，默认为中文简体+英文
        
    Returns:
        包含识别结果的字典
    """
    try:
        # 确保pytesseract库已安装
        if 'pytesseract' not in sys.modules:
            return {"success": False, "message": "缺少必要依赖: pytesseract"}
        
        # 检查图片文件是否存在
        if not os.path.exists(image_path):
            return {"success": False, "message": f"图片文件不存在: {image_path}"}
        
        # 打开图片并进行预处理
        image = Image.open(image_path)
        
        # 设置tesseract语言
        pytesseract.pytesseract.tesseract_cmd = r'tesseract'  # 默认命令，用户可能需要配置完整路径
        
        # 进行OCR识别
        start_time = time.time()
        text = pytesseract.image_to_string(image, lang=language)
        end_time = time.time()
        
        # 保存识别结果到txt文件
        output_file = os.path.join(OUTPUT_DIR, os.path.splitext(os.path.basename(image_path))[0] + "_ocr.txt")
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(text)
        
        return {
            "success": True,
            "message": f"OCR识别成功，用时{end_time-start_time:.2f}秒",
            "text": text,
            "output_file": output_file
        }
        
    except Exception as e:
        logger.error(f"OCR识别过程中出错: {str(e)}")
        return {"success": False, "message": f"OCR识别失败: {str(e)}"}

# 文档比较功能
@mcp.tool()
def compare_documents(doc1_path: str, doc2_path: str, output_format: str = "html") -> Dict[str, Any]:
    """
    比较两个文档的内容差异
    
    Args:
        doc1_path: 第一个文档的路径
        doc2_path: 第二个文档的路径
        output_format: 输出格式，支持'html'或'text'
        
    Returns:
        包含比较结果的字典
    """
    try:
        # 检查文件是否存在
        if not os.path.exists(doc1_path):
            return {"success": False, "message": f"文件不存在: {doc1_path}"}
        if not os.path.exists(doc2_path):
            return {"success": False, "message": f"文件不存在: {doc2_path}"}
        
        # 根据文件扩展名处理不同类型的文档
        ext1 = os.path.splitext(doc1_path)[1].lower()
        ext2 = os.path.splitext(doc2_path)[1].lower()
        
        # 提取文档文本内容
        text1 = extract_document_text(doc1_path)
        text2 = extract_document_text(doc2_path)
        
        if text1 is None or text2 is None:
            return {"success": False, "message": "无法提取文档文本内容"}
        
        # 按行分割文本
        lines1 = text1.splitlines()
        lines2 = text2.splitlines()
        
        # 使用difflib比较文本差异
        if output_format == 'html':
            diff = difflib.HtmlDiff()
            result = diff.make_file(lines1, lines2, os.path.basename(doc1_path), os.path.basename(doc2_path))
            output_file = os.path.join(OUTPUT_DIR, "document_diff.html")
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(result)
        else:  # text format
            diff = difflib.unified_diff(lines1, lines2, os.path.basename(doc1_path), os.path.basename(doc2_path))
            result = '\n'.join(list(diff))
            output_file = os.path.join(OUTPUT_DIR, "document_diff.txt")
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(result)
        
        return {
            "success": True,
            "message": "文档比较完成",
            "output_file": output_file,
            "difference_count": len([line for line in difflib.ndiff(lines1, lines2) if line.startswith('+ ') or line.startswith('- ')])
        }
        
    except Exception as e:
        logger.error(f"文档比较过程中出错: {str(e)}")
        return {"success": False, "message": f"文档比较失败: {str(e)}"}

# 文档翻译功能
@mcp.tool()
def translate_document(doc_path: str, target_language: str = "ZH", api_key: str = None) -> Dict[str, Any]:
    """
    将文档内容翻译成指定目标语言
    
    Args:
        doc_path: 文档路径
        target_language: 目标语言，默认为中文
        api_key: DeepL API密钥，可选
        
    Returns:
        包含翻译结果的字典
    """
    try:
        # 确保deepl库已安装
        if 'deepl' not in sys.modules:
            return {"success": False, "message": "缺少必要依赖: deepl"}
        
        # 检查文件是否存在
        if not os.path.exists(doc_path):
            return {"success": False, "message": f"文件不存在: {doc_path}"}
        
        # 提取文档文本
        text = extract_document_text(doc_path)
        if text is None:
            return {"success": False, "message": "无法提取文档文本内容"}
        
        # 初始化DeepL翻译器
        if api_key:
            translator = deepl.Translator(api_key)
        else:
            # 尝试从环境变量获取API密钥
            env_api_key = os.environ.get("DEEPL_API_KEY")
            if not env_api_key:
                return {"success": False, "message": "未提供DeepL API密钥"}
            translator = deepl.Translator(env_api_key)
        
        # 执行翻译
        start_time = time.time()
        result = translator.translate_text(text, target_lang=target_language)
        end_time = time.time()
        
        translated_text = result.text if hasattr(result, 'text') else str(result)
        
        # 保存翻译结果
        filename = os.path.splitext(os.path.basename(doc_path))[0]
        output_file = os.path.join(OUTPUT_DIR, f"{filename}_translated_{target_language}.txt")
        with open(output_file, "w", encoding="utf-8") as f:
            f.write(translated_text)
        
        return {
            "success": True,
            "message": f"文档翻译完成，用时{end_time-start_time:.2f}秒",
            "translated_text": translated_text,
            "output_file": output_file,
            "target_language": target_language
        }
        
    except Exception as e:
        logger.error(f"文档翻译过程中出错: {str(e)}")
        return {"success": False, "message": f"文档翻译失败: {str(e)}"}

# 文档加密功能
@mcp.tool()
def encrypt_document(doc_path: str, password: str) -> Dict[str, Any]:
    """
    加密文档
    
    Args:
        doc_path: 要加密的文档路径
        password: 加密密码
        
    Returns:
        包含加密结果的字典
    """
    try:
        # 确保cryptography库已安装
        if 'cryptography' not in sys.modules:
            return {"success": False, "message": "缺少必要依赖: cryptography"}
        
        # 检查文件是否存在
        if not os.path.exists(doc_path):
            return {"success": False, "message": f"文件不存在: {doc_path}"}
        
        # 读取文件内容
        with open(doc_path, 'rb') as file:
            file_data = file.read()
        
        # 生成密钥
        password_bytes = password.encode('utf-8')
        key = base64.urlsafe_b64encode(hashlib.sha256(password_bytes).digest())
        
        # 使用Fernet对称加密
        fernet = Fernet(key)
        encrypted_data = fernet.encrypt(file_data)
        
        # 保存加密文件
        filename = os.path.splitext(os.path.basename(doc_path))[0]
        output_file = os.path.join(OUTPUT_DIR, f"{filename}_encrypted.bin")
        with open(output_file, 'wb') as file:
            file.write(encrypted_data)
        
        return {
            "success": True,
            "message": "文档加密完成",
            "output_file": output_file
        }
        
    except Exception as e:
        logger.error(f"文档加密过程中出错: {str(e)}")
        return {"success": False, "message": f"文档加密失败: {str(e)}"}

# 文档解密功能
@mcp.tool()
def decrypt_document(encrypted_file_path: str, password: str, output_format: str = None) -> Dict[str, Any]:
    """
    解密文档
    
    Args:
        encrypted_file_path: 加密文件路径
        password: 解密密码
        output_format: 输出文件格式后缀，如果不提供则保持原格式
        
    Returns:
        包含解密结果的字典
    """
    try:
        # 确保cryptography库已安装
        if 'cryptography' not in sys.modules:
            return {"success": False, "message": "缺少必要依赖: cryptography"}
        
        # 检查文件是否存在
        if not os.path.exists(encrypted_file_path):
            return {"success": False, "message": f"文件不存在: {encrypted_file_path}"}
        
        # 读取加密文件
        with open(encrypted_file_path, 'rb') as file:
            encrypted_data = file.read()
        
        # 生成密钥
        password_bytes = password.encode('utf-8')
        key = base64.urlsafe_b64encode(hashlib.sha256(password_bytes).digest())
        
        # 解密数据
        try:
            fernet = Fernet(key)
            decrypted_data = fernet.decrypt(encrypted_data)
        except Exception as e:
            return {"success": False, "message": f"解密失败，密码可能不正确: {str(e)}"}
        
        # 保存解密文件
        output_ext = output_format if output_format else ".bin"
        if not output_ext.startswith('.'):
            output_ext = '.' + output_ext
            
        filename = os.path.splitext(os.path.basename(encrypted_file_path))[0]
        if filename.endswith("_encrypted"):
            filename = filename[:-10]  # 移除"_encrypted"后缀
            
        output_file = os.path.join(OUTPUT_DIR, f"{filename}_decrypted{output_ext}")
        with open(output_file, 'wb') as file:
            file.write(decrypted_data)
        
        return {
            "success": True,
            "message": "文档解密完成",
            "output_file": output_file
        }
        
    except Exception as e:
        logger.error(f"文档解密过程中出错: {str(e)}")
        return {"success": False, "message": f"文档解密失败: {str(e)}"}

# Excel数据导出到数据库功能
@mcp.tool()
def export_excel_to_database(excel_file: str, db_connection_string: str, table_name: str,
                          sheet_name: str = None, if_exists: str = "replace") -> Dict[str, Any]:
    """
    将Excel数据导出到数据库
    
    Args:
        excel_file: Excel文件路径
        db_connection_string: 数据库连接字符串
        table_name: 目标表名
        sheet_name: 工作表名称，默认为第一个工作表
        if_exists: 如果表已存在的处理方式，"replace"替换,"append"追加,"fail"报错
        
    Returns:
        包含操作结果的字典
    """
    try:
        # 确保pandas和sqlalchemy库已安装
        if 'pandas' not in sys.modules or 'sqlalchemy' not in sys.modules:
            return {"success": False, "message": "缺少必要依赖: pandas 或 sqlalchemy"}
        
        # 检查文件是否存在
        if not os.path.exists(excel_file):
            return {"success": False, "message": f"文件不存在: {excel_file}"}
        
        # 从Excel读取数据
        start_time = time.time()
        df = pd.read_excel(excel_file, sheet_name=sheet_name)
        
        # 创建数据库引擎
        engine = create_engine(db_connection_string)
        
        # 将数据导出到数据库
        df.to_sql(table_name, engine, if_exists=if_exists, index=False)
        end_time = time.time()
        
        # 获取行数和列数
        row_count = len(df)
        column_count = len(df.columns)
        
        return {
            "success": True,
            "message": f"数据成功导出到数据库表 {table_name}，用时{end_time-start_time:.2f}秒",
            "row_count": row_count,
            "column_count": column_count,
            "table_name": table_name
        }
        
    except Exception as e:
        logger.error(f"Excel数据导出到数据库过程中出错: {str(e)}")
        return {"success": False, "message": f"Excel数据导出失败: {str(e)}"}

# 数据库数据导入到Excel功能
@mcp.tool()
def import_database_to_excel(db_connection_string: str, query: str, output_file: str = None) -> Dict[str, Any]:
    """
    从数据库导入数据到Excel文件
    
    Args:
        db_connection_string: 数据库连接字符串
        query: SQL查询语句
        output_file: 输出Excel文件名，默认为query_result.xlsx
        
    Returns:
        包含操作结果的字典
    """
    try:
        # 确保pandas和sqlalchemy库已安装
        if 'pandas' not in sys.modules or 'sqlalchemy' not in sys.modules:
            return {"success": False, "message": "缺少必要依赖: pandas 或 sqlalchemy"}
        
        # 创建数据库引擎
        engine = create_engine(db_connection_string)
        
        # 执行SQL查询
        start_time = time.time()
        df = pd.read_sql(query, engine)
        
        # 设置输出文件名
        if not output_file:
            output_file = os.path.join(OUTPUT_DIR, "query_result.xlsx")
        elif not os.path.isabs(output_file):
            output_file = os.path.join(OUTPUT_DIR, output_file)
            
        # 导出数据到Excel
        df.to_excel(output_file, index=False)
        end_time = time.time()
        
        # 获取行数和列数
        row_count = len(df)
        column_count = len(df.columns)
        
        return {
            "success": True,
            "message": f"数据库查询结果已成功导出到Excel，用时{end_time-start_time:.2f}秒",
            "output_file": output_file,
            "row_count": row_count,
            "column_count": column_count
        }
        
    except Exception as e:
        logger.error(f"数据库数据导入到Excel过程中出错: {str(e)}")
        return {"success": False, "message": f"数据库数据导入失败: {str(e)}"}

# 批量生成文档功能
@mcp.tool()
def batch_create_documents(template_path: str, output_prefix: str, count: int, 
                        data_mapping: Dict[str, List[str]] = None) -> Dict[str, Any]:
    """
    基于模板批量生成文档
    
    Args:
        template_path: 模板文档路径
        output_prefix: 输出文件前缀
        count: 生成的文档数量
        data_mapping: 数据映射字典，键为占位符名称，值为数据列表
        
    Returns:
        包含操作结果的字典
    """
    try:
        # 检查模板文件是否存在
        if not os.path.exists(template_path):
            return {"success": False, "message": f"模板文件不存在: {template_path}"}
        
        # 确保数据足够生成所需数量的文档
        if data_mapping:
            min_data_count = min([len(values) for values in data_mapping.values()])
            if min_data_count < count:
                return {"success": False, "message": f"数据不足，最少需要{count}条记录，但只有{min_data_count}条"}
        
        # 获取文件扩展名
        _, ext = os.path.splitext(template_path)
        
        # 创建输出目录
        if not os.path.exists(OUTPUT_DIR):
            os.makedirs(OUTPUT_DIR)
        
        # 批量生成文档
        start_time = time.time()
        output_files = []
        
        for i in range(count):
            # 创建输出文件路径
            output_file = os.path.join(OUTPUT_DIR, f"{output_prefix}_{i+1}{ext}")
            
            # 复制模板文件
            shutil.copy2(template_path, output_file)
            
            # 如果有数据映射，替换占位符
            if data_mapping:
                replace_placeholders(output_file, data_mapping, i)
            
            output_files.append(output_file)
        
        end_time = time.time()
        
        return {
            "success": True,
            "message": f"成功生成{count}个文档，用时{end_time-start_time:.2f}秒",
            "output_files": output_files
        }
        
    except Exception as e:
        logger.error(f"批量生成文档过程中出错: {str(e)}")
        return {"success": False, "message": f"批量生成文档失败: {str(e)}"}

@mcp.tool()
def batch_process_documents(files: List[str], operation: str, params: Dict[str, Any] = None,
                         max_workers: int = 4) -> Dict[str, Any]:
    """
    并行批量处理多个文档
    
    Args:
        files: 要处理的文件路径列表
        operation: 要执行的操作名称，如'translate_document'
        params: 操作参数
        max_workers: 最大并行工作线程数
        
    Returns:
        包含操作结果的字典
    """
    try:
        # 获取所有工具函数
        operations_mapping = {
            "translate_document": translate_document,
            "encrypt_document": encrypt_document,
            "decrypt_document": decrypt_document,
            "ocr_recognize_text": ocr_recognize_text,
            "compare_documents": compare_documents,
            "export_excel_to_database": export_excel_to_database,
            "import_database_to_excel": import_database_to_excel,
            "general_file_operations": general_file_operations
        }
        
        # 检查操作是否支持
        if operation not in operations_mapping:
            return {"success": False, "message": f"不支持的操作: {operation}"}
        
        # 获取操作函数
        op_func = operations_mapping[operation]
        if params is None:
            params = {}
        
        # 初始化结果列表
        results = []
        success_count = 0
        failed_files = []
        
        # 定义并行处理的任务函数
        def process_file(file_path):
            try:
                # 复制参数并添加文件路径
                file_params = params.copy()
                
                # 根据不同操作类型，设置不同的文件参数名称
                param_name_mapping = {
                    "translate_document": "doc_path",
                    "encrypt_document": "doc_path",
                    "decrypt_document": "encrypted_file_path",
                    "ocr_recognize_text": "image_path",
                    "compare_documents": "doc1_path"  # 注意: 这种情况可能需要特殊处理
                }
                
                param_name = param_name_mapping.get(operation, "file_path")
                file_params[param_name] = file_path
                
                # 调用操作函数
                result = op_func(**file_params)
                
                # 添加文件路径到结果
                result["file_path"] = file_path
                return result
            except Exception as e:
                logger.error(f"处理文件 {file_path} 时出错: {str(e)}")
                return {"success": False, "message": str(e), "file_path": file_path}
        
        # 使用线程池并行处理文件
        start_time = time.time()
        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_file = {executor.submit(process_file, file): file for file in files}
            
            for future in concurrent.futures.as_completed(future_to_file):
                file = future_to_file[future]
                try:
                    result = future.result()
                    results.append(result)
                    
                    if result.get("success", False):
                        success_count += 1
                    else:
                        failed_files.append(file)
                        
                except Exception as e:
                    logger.error(f"获取任务结果时出错 {file}: {str(e)}")
                    results.append({"success": False, "message": str(e), "file_path": file})
                    failed_files.append(file)
        
        end_time = time.time()
        
        return {
            "success": True,
            "message": f"批量处理完成，成功: {success_count}/{len(files)}，用时: {end_time-start_time:.2f}秒",
            "results": results,
            "success_count": success_count,
            "failed_files": failed_files
        }
        
    except Exception as e:
        logger.error(f"批量处理文档过程中出错: {str(e)}")
        return {"success": False, "message": f"批量处理文档失败: {str(e)}"}

# 文件目录操作功能
@mcp.tool()
def general_file_operations(operation: str, source_path: str, 
                         target_path: str = None, recursive: bool = False) -> Dict[str, Any]:
    """
    通用文件和目录操作
    
    Args:
        operation: 操作类型，"copy"复制, "move"移动, "delete"删除, "list"列表
        source_path: 源文件或目录路径
        target_path: 目标路径，用于复制和移动操作
        recursive: 是否递归处理子目录
        
    Returns:
        包含操作结果的字典
    """
    try:
        # 检查源路径是否存在
        if not os.path.exists(source_path):
            return {"success": False, "message": f"源路径不存在: {source_path}"}
        
        # 根据操作类型执行不同的操作
        if operation == "copy":
            if not target_path:
                return {"success": False, "message": "复制操作需要提供目标路径"}
                
            if os.path.isdir(source_path):
                # 复制目录
                if recursive:
                    shutil.copytree(source_path, target_path, dirs_exist_ok=True)
                else:
                    # 仅复制目录中的文件，不包括子目录
                    if not os.path.exists(target_path):
                        os.makedirs(target_path)
                    for item in os.listdir(source_path):
                        item_path = os.path.join(source_path, item)
                        if os.path.isfile(item_path):
                            shutil.copy2(item_path, target_path)
                return {"success": True, "message": f"成功将{source_path}复制到{target_path}"}
            else:
                # 复制文件
                if os.path.isdir(target_path):
                    # 如果目标是目录，则复制到该目录中
                    shutil.copy2(source_path, target_path)
                    target_file = os.path.join(target_path, os.path.basename(source_path))
                else:
                    # 如果目标是文件路径，则直接复制
                    shutil.copy2(source_path, target_path)
                    target_file = target_path
                return {"success": True, "message": f"成功将{source_path}复制到{target_file}"}
                
        elif operation == "move":
            if not target_path:
                return {"success": False, "message": "移动操作需要提供目标路径"}
                
            shutil.move(source_path, target_path)
            return {"success": True, "message": f"成功将{source_path}移动到{target_path}"}
            
        elif operation == "delete":
            if os.path.isdir(source_path):
                if recursive:
                    shutil.rmtree(source_path)
                else:
                    # 只删除目录中的文件，保留目录结构
                    for item in os.listdir(source_path):
                        item_path = os.path.join(source_path, item)
                        if os.path.isfile(item_path):
                            os.remove(item_path)
                return {"success": True, "message": f"成功删除目录{source_path}"}
            else:
                os.remove(source_path)
                return {"success": True, "message": f"成功删除文件{source_path}"}
                
        elif operation == "list":
            if os.path.isdir(source_path):
                # 列出目录内容
                file_list = []
                
                if recursive:
                    # 递归列出所有文件和目录
                    for root, dirs, files in os.walk(source_path):
                        for file in files:
                            full_path = os.path.join(root, file)
                            relative_path = os.path.relpath(full_path, source_path)
                            file_list.append({
                                "name": file,
                                "path": relative_path,
                                "size": os.path.getsize(full_path),
                                "is_dir": False,
                                "modified": os.path.getmtime(full_path)
                            })
                        for dir_name in dirs:
                            full_path = os.path.join(root, dir_name)
                            relative_path = os.path.relpath(full_path, source_path)
                            file_list.append({
                                "name": dir_name,
                                "path": relative_path,
                                "is_dir": True,
                                "modified": os.path.getmtime(full_path)
                            })
                else:
                    # 仅列出当前目录内容
                    for item in os.listdir(source_path):
                        item_path = os.path.join(source_path, item)
                        is_dir = os.path.isdir(item_path)
                        info = {
                            "name": item,
                            "path": item,
                            "is_dir": is_dir,
                            "modified": os.path.getmtime(item_path)
                        }
                        if not is_dir:
                            info["size"] = os.path.getsize(item_path)
                        file_list.append(info)
                        
                return {
                    "success": True,
                    "message": f"成功列出{source_path}中的{len(file_list)}个项目",
                    "items": file_list
                }
            else:
                # 返回单个文件的信息
                file_info = {
                    "name": os.path.basename(source_path),
                    "path": source_path,
                    "size": os.path.getsize(source_path),
                    "is_dir": False,
                    "modified": os.path.getmtime(source_path)
                }
                return {
                    "success": True,
                    "message": "成功获取文件信息",
                    "items": [file_info]
                }
        else:
            return {"success": False, "message": f"不支持的操作: {operation}"}
            
    except Exception as e:
        logger.error(f"文件操作过程中出错: {str(e)}")
        return {"success": False, "message": f"文件操作失败: {str(e)}"}

# 启动MCP服务器
if __name__ == "__main__":
    logger.info("正在启动高级Office功能服务器...")
    mcp.run()
